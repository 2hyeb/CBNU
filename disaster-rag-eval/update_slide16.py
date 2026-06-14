"""Update PPT slide 16 (index 15) with actual q_type RAGAS results."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

PPT_PATH = r"D:\충북대\지능화캡스톤\지능화캡스톤프로젝트_최종_수정_2025254022_이혜빈_2.pptx"

# q_type results: rows = [절차형, 안전판단형, 복합추론형, 수치법령형]
# cols = [EXAONE/rag, Qwen3/rag, LLaMA/rag, EXAONE/ft_only, Qwen3/ft_only, LLaMA/ft_only]
DATA = [
    ["0.833", "0.901", "0.136", "평가불가", "평가불가", "평가불가"],
    ["0.751", "0.773", "0.000", "평가불가", "평가불가", "평가불가"],
    ["0.647", "0.702", "N/A",  "평가불가", "평가불가", "평가불가"],
    ["0.756", "0.742", "0.000", "평가불가", "평가불가", "평가불가"],
]

COLOR_VALUE  = RGBColor(0, 0, 0)        # normal value
COLOR_NA     = RGBColor(150, 150, 150)  # N/A
COLOR_FAIL   = RGBColor(192, 0, 0)      # 평가불가


def set_cell_text(cell, text, bold=False, italic=False, color=None, size=9, align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # Remove all existing runs cleanly
    for old_run in p.runs:
        p._p.remove(old_run._r)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


prs = Presentation(PPT_PATH)
slide = prs.slides[15]  # slide 16

# Identify tables on slide
tables = [sh for sh in slide.shapes if sh.has_table]
print(f"Tables found on slide 16: {len(tables)}")
for i, t in enumerate(tables):
    tbl = t.table
    print(f"  Table {i}: {len(tbl.rows)} rows × {len(tbl.columns)} cols")

if len(tables) < 2:
    print("ERROR: expected at least 2 tables")
    exit(1)

# Main results table: 5 rows × 7 cols
main_tbl = None
hyp_tbl = None
for t in tables:
    tbl = t.table
    nrows = len(tbl.rows)
    ncols = len(tbl.columns)
    print(f"    checking {nrows}×{ncols}")
    if nrows == 5 and ncols == 7:
        main_tbl = tbl
    elif nrows == 4 and ncols == 4:
        hyp_tbl = tbl

if main_tbl is None:
    print("ERROR: main results table (5×7) not found")
    exit(1)
if hyp_tbl is None:
    print("ERROR: hypothesis table (4×4) not found")
    exit(1)

print("\nUpdating main results table...")
for ri, row_data in enumerate(DATA):
    q_type_names = ["절차형", "안전판단형", "복합추론형", "수치법령형"]
    print(f"  Row {ri+1} ({q_type_names[ri]}): {row_data}")
    for ci, val in enumerate(row_data):
        cell = main_tbl.cell(ri + 1, ci + 1)  # skip header row and q_type column
        if val == "평가불가":
            set_cell_text(cell, val, bold=False, italic=True, color=COLOR_FAIL)
        elif val == "N/A":
            set_cell_text(cell, val, bold=False, italic=True, color=COLOR_NA)
        else:
            set_cell_text(cell, val, bold=False, italic=False, color=COLOR_VALUE)

print("\nUpdating hypothesis table...")
# H1 row (ri=1), col 2 = "현재 상태"
# H2 row (ri=2), col 2 = "현재 상태"
h1_status = "부분 지지"   # 수치법령형: EXAONE/rag 0.756, Qwen3/rag 0.742 > ft_only(평가불가)
h2_status = "검증 불가"   # ft_only 평가불가로 FT 우세 확인 불가
h1_evidence = "수치법령형 RAG: EXAONE 0.756 / Qwen3 0.742 — RAG 우수 확인. ft_only 평가불가(RAGAS 파싱 오류)로 직접 비교 제한"
h2_evidence = "ft_only 모델 전체 answer_relevancy ≈ 0.0 (RAGAS 평가 불가 — LLM 심판 응답 파싱 실패). FT 우세 가설 수치 검증 불가"

set_cell_text(hyp_tbl.cell(1, 2), h1_status, color=RGBColor(0, 112, 0), bold=True)
set_cell_text(hyp_tbl.cell(2, 2), h2_status, color=RGBColor(192, 0, 0), bold=True)
set_cell_text(hyp_tbl.cell(1, 3), h1_evidence, align=PP_ALIGN.LEFT, size=8)
set_cell_text(hyp_tbl.cell(2, 3), h2_evidence, align=PP_ALIGN.LEFT, size=8)

print("\nUpdating footnote textbox...")
NEW_FOOTNOTE = (
    "* Faithfulness (RAG 조건) / Answer Relevancy (ft_only 조건) 기준  |  "
    "평가불가: ft_only 모델 RAGAS answer_relevancy ≈ 0.0 — LLM 심판(Qwen3-32B) 응답 파싱 실패(OutputParserException)로 자동 평가 불가  |  "
    "N/A: LLaMA/rag 복합추론형 faithfulness NaN"
)

# Find footnote textbox (contains the old footnote text about "진행 중")
footnote_updated = False
for sh in slide.shapes:
    if sh.has_text_frame:
        full_text = sh.text_frame.text
        if "Faithfulness" in full_text and "진행 중" in full_text:
            tf = sh.text_frame
            # Clear all paragraphs and set single run
            for i, para in enumerate(tf.paragraphs):
                for run in para.runs:
                    para._p.remove(run._r)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = NEW_FOOTNOTE
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(128, 128, 128)
            footnote_updated = True
            print(f"  Updated footnote textbox")
            break

if not footnote_updated:
    print("  WARNING: footnote textbox not found, skipping")

prs.save(PPT_PATH)
print(f"\nSaved: {PPT_PATH}")
print("Slide 16 update complete.")
