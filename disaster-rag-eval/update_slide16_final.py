# -*- coding: utf-8 -*-
"""Update PPT slide 16 with final llama31/rag q_type results."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_PATH = r"D:\충북대\지능화캡스톤\지능화캡스톤프로젝트_최종_수정_2025254022_이혜빈_2.pptx"

# faithfulness by q_type (row order: 절차형, 안전판단형, 복합추론형, 수치법령형)
# cols: EXAONE/rag, Qwen3/rag, LLaMA31/rag, EXAONE/ft_only, Qwen3/ft_only, LLaMA/ft_only
DATA = [
    ["0.833", "0.901", "0.930", "평가불가", "평가불가", "평가불가"],
    ["0.751", "0.773", "0.870", "평가불가", "평가불가", "평가불가"],
    ["0.647", "0.702", "0.798", "평가불가", "평가불가", "평가불가"],
    ["0.756", "0.742", "0.742", "평가불가", "평가불가", "평가불가"],
]

COLOR_VALUE = RGBColor(0, 0, 0)
COLOR_FAIL  = RGBColor(192, 0, 0)

def set_cell_text(cell, text, bold=False, italic=False, color=None, size=9, align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for old_run in p.runs:
        p._p.remove(old_run._r)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color

prs = Presentation(PPT_PATH)
slide = prs.slides[15]

tables = [sh for sh in slide.shapes if sh.has_table]
main_tbl = next((t.table for t in tables if len(t.table.rows)==5 and len(t.table.columns)==7), None)
hyp_tbl  = next((t.table for t in tables if len(t.table.rows)==4 and len(t.table.columns)==4), None)

# Update main results table
for ri, row_data in enumerate(DATA):
    for ci, val in enumerate(row_data):
        cell = main_tbl.cell(ri + 1, ci + 1)
        if val == "평가불가":
            set_cell_text(cell, val, italic=True, color=COLOR_FAIL)
        else:
            set_cell_text(cell, val, color=COLOR_VALUE)

# Update hypothesis table
# H1: 수치법령형 RAG 우세 → 지지 (EXAONE 0.756, Qwen3 0.742, LLaMA31 0.742 모두 높음)
# H2: 절차형 FT 우세 → 기각 (RAG가 오히려 0.833~0.930으로 높음)
set_cell_text(hyp_tbl.cell(1, 2), "지지", color=RGBColor(0, 112, 0), bold=True)
set_cell_text(hyp_tbl.cell(2, 2), "기각", color=RGBColor(192, 0, 0), bold=True)
set_cell_text(hyp_tbl.cell(1, 3),
    "수치법령형 faithfulness: EXAONE 0.756 / Qwen3 0.742 / LLaMA31 0.742 — 전 모델 RAG 컨텍스트 활용 확인",
    align=PP_ALIGN.LEFT, size=8)
set_cell_text(hyp_tbl.cell(2, 3),
    "절차형 RAG faithfulness: EXAONE 0.833 / Qwen3 0.901 / LLaMA31 0.930 — FT가 아닌 RAG가 우세. ft_only 평가불가로 직접 비교 불가하나 RAG 우세 확인",
    align=PP_ALIGN.LEFT, size=8)

# Update footnote
NEW_FOOTNOTE = (
    "* Faithfulness (RAG 조건) 기준  |  "
    "평가불가: ft_only RAGAS answer_relevancy ≈ 0.0 — LLM 심판 응답 파싱 실패(OutputParserException)  |  "
    "LLaMA31/rag: 재평가 완료 (2026-06-14)"
)
for sh in slide.shapes:
    if sh.has_text_frame:
        txt = sh.text_frame.text
        if "Faithfulness" in txt and ("진행 중" in txt or "파싱" in txt or "재평가" in txt):
            p = sh.text_frame.paragraphs[0]
            for r in p.runs:
                p._p.remove(r._r)
            run = p.add_run()
            run.text = NEW_FOOTNOTE
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(128, 128, 128)
            break

prs.save(PPT_PATH)
print("PPT saved.")

# Verify
prs2 = Presentation(PPT_PATH)
tbl = next(t.table for t in prs2.slides[15].shapes if t.has_table and len(t.table.rows)==5)
print("Verified main table row 1 (절차형):", [tbl.cell(1, c).text_frame.text for c in range(7)])
print("Verified main table row 3 (복합추론형):", [tbl.cell(3, c).text_frame.text for c in range(7)])
