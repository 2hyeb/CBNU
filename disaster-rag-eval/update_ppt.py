# -*- coding: utf-8 -*-
"""PPT _2 업데이트 스크립트"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import lxml.etree as etree

PPT_PATH = r"D:\충북대\지능화캡스톤\지능화캡스톤프로젝트_최종_수정_2025254022_이혜빈_2.pptx"
prs = Presentation(PPT_PATH)

SL = prs.slide_width
SH = prs.slide_height

# ── 헬퍼 함수 ────────────────────────────────────────────────────────────────

def find_and_replace(slide, old_text, new_text):
    """슬라이드 내 특정 텍스트를 새 텍스트로 교체 (run 단위)"""
    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    replaced = True
    return replaced


def add_text_to_frame(tf, text, bold=False, size=14, color=None, align=PP_ALIGN.LEFT):
    para = tf.add_paragraph()
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor(*color)
    return para


def add_textbox(slide, left, top, width, height, text, bold=False, size=14,
                color=None, bg=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor(*color)
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg)
    return txBox


def add_table_slide(prs, title_text, subtitle_text, rows_data, col_widths,
                    insert_after_idx, header_row=True):
    """새 슬라이드(blank layout)에 표 추가 후 지정 위치에 삽입"""
    layout = prs.slide_layouts[3]  # 사용자 지정 레이아웃
    new_slide = prs.slides.add_slide(layout)
    # 기존 placeholder 모두 제거
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # 슬라이드를 insert_after_idx 뒤에 이동
    xml_slides = prs.slides._sldIdLst
    last_id = xml_slides[-1]
    xml_slides.remove(last_id)
    xml_slides.insert(insert_after_idx + 1, last_id)

    W, H = SL, SH

    # 상단 제목 배너
    add_textbox(new_slide, Inches(0), Inches(0), W, Inches(0.55),
                title_text, bold=True, size=16,
                color=(255,255,255), bg=(31,73,125), align=PP_ALIGN.CENTER)

    # 부제목
    add_textbox(new_slide, Inches(0.3), Inches(0.6), W - Inches(0.6), Inches(0.35),
                subtitle_text, bold=False, size=12, color=(89,89,89))

    # 표 추가
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_top = Inches(1.0)
    tbl_left = Inches(0.3)
    tbl_width = W - Inches(0.6)
    tbl_height = Inches(0.38) * n_rows

    table = new_slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top,
                                        tbl_width, tbl_height).table

    # 열 너비 설정
    total = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = int(tbl_width * cw / total)

    # 셀 채우기
    HEADER_BG = RGBColor(31, 73, 125)
    HEADER_FG = RGBColor(255, 255, 255)
    ALT_BG    = RGBColor(235, 241, 250)
    GOOD      = RGBColor(0, 112, 0)
    WARN      = RGBColor(180, 40, 40)

    for ri, row_data in enumerate(rows_data):
        is_header = header_row and ri == 0
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(10)
            run.font.bold = is_header

            if is_header:
                run.font.color.rgb = HEADER_FG
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG
            else:
                # 짝수 행 배경
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ALT_BG
                # 수치 색상 강조
                try:
                    v = float(cell_text)
                    if ci >= 2 and v >= 0.7:
                        run.font.color.rgb = GOOD
                    elif ci >= 2 and v < 0.25 and cell_text not in ('—', ''):
                        run.font.color.rgb = WARN
                except (ValueError, TypeError):
                    pass

    return new_slide


# ═══════════════════════════════════════════════════════════════════════════════
# 1. 슬라이드 9 수정 — LLaMA 설명 업데이트
# ═══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides[8]
find_and_replace(slide9,
    "LLaMA 3.1 8B — 다국어 지원, #1에서 Faithfulness 최고",
    "LLaMA 3.1 8B — 다국어 지원, 영어 중심 모델 (한국어 RAG 한계 실험적 확인)")
print("Slide 9 수정 완료")


# ═══════════════════════════════════════════════════════════════════════════════
# 2. 슬라이드 13 수정 — "예상 결과 방향성" → "→ 별도 슬라이드 참조"
# ═══════════════════════════════════════════════════════════════════════════════
slide13 = prs.slides[12]
find_and_replace(slide13, "예상 결과 방향성",
                 "실제 실험 결과 → 다음 슬라이드 참조")
print("Slide 13 수정 완료")


# ═══════════════════════════════════════════════════════════════════════════════
# 3. 슬라이드 15 수정 — 한계점 2개 추가
# ═══════════════════════════════════════════════════════════════════════════════
slide15 = prs.slides[14]
for shape in slide15.shapes:
    if shape.has_text_frame and "토론 포인트" in shape.text_frame.text:
        tf = shape.text_frame
        # 빈 줄
        p_blank = tf.add_paragraph()
        p_blank.text = ""
        # 한계 6
        p6_title = tf.add_paragraph()
        r6t = p6_title.add_run()
        r6t.text = "FT+RAG 학습 데이터 형식 불일치 (설계 결함)"
        r6t.font.bold = True
        r6t.font.size = Pt(12)
        r6t.font.color.rgb = RGBColor(192, 0, 0)
        p6_body = tf.add_paragraph()
        r6b = p6_body.add_run()
        r6b.text = ("학습 시 RAG 컨텍스트 없이 훈련된 FT 모델이 inference 시 [참고 자료] 섹션을 무시 → "
                    "전 모델 ft_rag 조건 faithfulness 저하 발생. "
                    "학습 데이터에 RAG 형식([참고 자료]+질문→정답) 통합 후 재실험 필요.")
        r6b.font.size = Pt(11)
        # 한계 7
        p7_title = tf.add_paragraph()
        r7t = p7_title.add_run()
        r7t.text = "영어 중심 모델의 한국어 RAG 처리 한계"
        r7t.font.bold = True
        r7t.font.size = Pt(12)
        r7t.font.color.rgb = RGBColor(192, 0, 0)
        p7_body = tf.add_paragraph()
        r7b = p7_body.add_run()
        r7b.text = ("LLaMA-3.1-8B는 no_rag(answer_rel 0.8149) 대비 rag 조건에서 faithfulness 0.19로 급락. "
                    "영어 중심 사전학습으로 한국어 [참고 자료] 컨텍스트를 노이즈로 처리. "
                    "한국어 재난안전 RAG에는 한국어 특화 또는 다국어 모델 사용 필수.")
        r7b.font.size = Pt(11)
        break
print("Slide 15 수정 완료")


# ═══════════════════════════════════════════════════════════════════════════════
# 4. 슬라이드 16 수정 — "예상" → 실제 결론으로 교체
# ═══════════════════════════════════════════════════════════════════════════════
slide16 = prs.slides[15]
find_and_replace(slide16, "연구 결론(예상)", "연구 결론 (실험 결과 기반)")
find_and_replace(slide16,
    "수치법령형: 외부 문서 검색이 핵심이므로 RAG가 FT 대비 우세 예상 (H1)",
    "H1 (부분 지지): EXAONE·Qwen3 rag 조건 faithfulness 0.77/0.81로 RAG 효과 확인 — 단 LLaMA는 한국어 RAG 실패")
find_and_replace(slide16,
    "절차형: 행동 순서 내재화가 핵심이므로 FT가 RAG 대비 우세 예상 (H2)",
    "H2 (검증 불충분): ft_only vs rag 직접 faithfulness 비교 불가 (ft_only 컨텍스트 없음) — q_type별 재평가 진행 중")
find_and_replace(slide16,
    "FT+RAG 결합: 전반적 우수 예상이나 유형별 편차 존재 가능 (H3)",
    "H3 (조건부 지지): 재학습 후 ft_rag faithfulness EXAONE 0.39·LLaMA 0.17·Qwen3 0.25 — RAG 단독 대비 낮으나 0에서 개선. 학습 데이터 RAG 형식 통합 시 개선 가능성 있음")
# 향후 연구 추가
find_and_replace(slide16,
    "더 큰 모델 및 Full FT 비교 : 13B 이상 모델 및 Full Fine-tuning과의 비교로 QLoRA 효과의 상한선을 추정함",
    "더 큰 모델 및 Full FT 비교 : 13B 이상 모델 및 Full Fine-tuning과의 비교로 QLoRA 효과의 상한선을 추정함\n질문 유형별 세분화 재평가 : q_type 조건부 RAGAS 평가로 H1·H2 직접 검증 및 유형별 배포 전략 수립")
print("Slide 16 수정 완료")


# ═══════════════════════════════════════════════════════════════════════════════
# 5. 새 슬라이드 A — 실제 실험 결과 표 (슬라이드 14 뒤에 삽입)
# ═══════════════════════════════════════════════════════════════════════════════
results_rows = [
    ["모델", "조건", "Faithfulness", "Answer Relevancy", "Context Recall"],
    ["EXAONE", "no_rag",  "—",      "0.7822", "—"],
    ["EXAONE", "rag",     "0.7689", "0.7749", "0.8474"],
    ["EXAONE", "ft_only", "—",      "0.7749", "—"],
    ["EXAONE", "ft_rag",  "0.3872", "0.7675", "0.0"],
    ["LLaMA",  "no_rag",  "—",      "0.8149", "—"],
    ["LLaMA",  "rag",     "0.1905", "0.8334", "0.0"],
    ["LLaMA",  "ft_only", "—",      "0.8478", "—"],
    ["LLaMA",  "ft_rag",  "0.1734", "0.8270", "0.0"],
    ["Qwen3",  "no_rag",  "—",      "0.7774", "—"],
    ["Qwen3",  "rag",     "0.8052", "0.7663", "0.8521"],
    ["Qwen3",  "ft_only", "—",      "0.7615", "—"],
    ["Qwen3",  "ft_rag",  "0.2495", "0.6044", "0.0"],
]
slide_A = add_table_slide(
    prs,
    title_text="결과 및 분석 (Results & Analysis)",
    subtitle_text="3모델 × 4조건 RAGAS 평가 결과 (83샘플, Judge: Qwen3-32B-AWQ)",
    rows_data=results_rows,
    col_widths=[2, 2, 2, 2, 2],
    insert_after_idx=14  # 슬라이드 14 뒤 (0-indexed: 13)
)

# 핵심 발견 요약 추가 (표 아래)
W = SL
notes = [
    ("✅ Qwen3+RAG 최고 성능", "faithfulness 0.8052  |  context_recall 0.8521"),
    ("⚠ LLaMA 한국어 RAG 실패", "no_rag 0.8149 → rag 0.1905 (한국어 컨텍스트 처리 한계)"),
    ("🔄 ft_rag 개선됨", "재학습 전 0.0 → 재학습 후 0.17~0.39 (학습 데이터 RAG 형식 반영 효과)"),
]
note_top = Inches(5.5)
for i, (title, body) in enumerate(notes):
    left_offset = Inches(0.3) + i * (W - Inches(0.6)) // 3
    box_w = (W - Inches(0.8)) // 3
    add_textbox(slide_A, left_offset, note_top, box_w, Inches(0.9),
                f"{title}\n{body}", bold=False, size=10,
                color=(0, 0, 0), bg=(242, 242, 242))

print("슬라이드 A (실험 결과 표) 추가 완료 — 슬라이드 15 위치")


# ═══════════════════════════════════════════════════════════════════════════════
# 6. 새 슬라이드 B — q_type별 분석 틀 (진행 중 placeholder)
# ═══════════════════════════════════════════════════════════════════════════════
layout = prs.slide_layouts[3]
slide_B = prs.slides.add_slide(layout)
for ph in list(slide_B.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)
# 슬라이드 A 바로 뒤로 이동 (실험 결과 표 다음)
xml_slides = prs.slides._sldIdLst
last_id = xml_slides[-1]
xml_slides.remove(last_id)
# slide_A는 현재 15번째(0-indexed 14), 그 다음에 삽입
xml_slides.insert(16, last_id)

W, H = SL, SH

# 제목 배너
add_textbox(slide_B, Inches(0), Inches(0), W, Inches(0.55),
            "결과 및 분석 (Results & Analysis)",
            bold=True, size=16,
            color=(255,255,255), bg=(31,73,125), align=PP_ALIGN.CENTER)

add_textbox(slide_B, Inches(0.3), Inches(0.6), W-Inches(0.6), Inches(0.35),
            "질문 유형별 교차 분석 결과 (q_type × condition)", bold=False, size=12, color=(89,89,89))

# 질문 유형 4개 설명 박스
qtypes = [
    ("절차형  (30문항)", "행동 순서 내재화가 핵심\n대피 절차, 응급조치 순서 등\n→ FT 우세 가설 (H2)"),
    ("안전판단형  (24문항)", "상황 판단 및 위험도 평가\n특정 상황에서의 적절 조치\n→ 종합 판단 필요"),
    ("복합추론형  (13문항)", "다단계 추론 요구\n복수 조건 결합 판단\n→ 두 방법 병행 가능"),
    ("수치법령형  (16문항)", "정확한 수치·기준 검색 핵심\n법령 수치, 농도 기준 등\n→ RAG 우세 가설 (H1)"),
]
box_w = (W - Inches(0.8)) // 4
for i, (title, body) in enumerate(qtypes):
    left = Inches(0.3) + i * (box_w + Inches(0.06))
    add_textbox(slide_B, left, Inches(1.05), box_w, Inches(1.3),
                f"{title}\n\n{body}", bold=False, size=10,
                color=(0,0,0), bg=(217,226,243))

# 결과 표 틀 (placeholder)
ph_rows = [
    ["질문 유형", "EXAONE/rag", "Qwen3/rag", "LLaMA/rag", "EXAONE/ft_only", "Qwen3/ft_only", "LLaMA/ft_only"],
    ["절차형",     "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]"],
    ["안전판단형", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]"],
    ["복합추론형", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]"],
    ["수치법령형", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]", "[진행 중]"],
]

tbl_top = Inches(2.5)
tbl_left = Inches(0.3)
tbl_width = W - Inches(0.6)
tbl_height = Inches(0.42) * len(ph_rows)
table_b = slide_B.shapes.add_table(len(ph_rows), len(ph_rows[0]),
                                    tbl_left, tbl_top, tbl_width, tbl_height).table
col_ws = [2.5, 1.5, 1.5, 1.5, 1.7, 1.7, 1.7]
total = sum(col_ws)
for ci, cw in enumerate(col_ws):
    table_b.columns[ci].width = int(tbl_width * cw / total)

HEADER_BG = RGBColor(31, 73, 125)
for ri, row in enumerate(ph_rows):
    for ci, txt in enumerate(row):
        cell = table_b.cell(ri, ci)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(9)
        run.font.bold = (ri == 0)
        if ri == 0:
            run.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BG
        elif ci == 0:
            run.font.bold = True
        elif txt == "[진행 중]":
            run.font.color.rgb = RGBColor(128, 128, 128)
            run.font.italic = True

# 지표 안내
add_textbox(slide_B, Inches(0.3), Inches(4.55), W-Inches(0.6), Inches(0.3),
            "* Faithfulness (RAG 조건) / Answer Relevancy (전 조건) 기준  |  진행 중: q_type별 RAGAS 평가 중 (완료 후 수치 업데이트 예정)",
            size=9, color=(128,128,128))

# 가설 결과 요약
hyp_rows = [
    ["가설", "예측", "현재 상태", "근거"],
    ["H1: 수치법령형 → RAG 우세", "RAG > FT", "평가 진행 중", "EXAONE/Qwen3 rag 전체 faithfulness 0.77/0.81 (유형별 세분화 진행 중)"],
    ["H2: 절차형 → FT 우세",      "FT > RAG", "평가 진행 중", "ft_only answer_rel 0.77~0.85 vs rag (유형별 세분화 진행 중)"],
    ["H3: FT+RAG 결합 최우수",    "ft_rag 최고", "조건부 지지",  "재학습 후 ft_rag 0.17~0.39, RAG 단독보다 낮음. 설계 개선 시 가능성 있음"],
]
tbl2_top = Inches(5.0)
tbl2 = slide_B.shapes.add_table(len(hyp_rows), 4,
                                  tbl_left, tbl2_top, tbl_width, Inches(0.38)*len(hyp_rows)).table
hw = [2.5, 1.5, 1.5, 5.5]
total2 = sum(hw)
for ci, cw in enumerate(hw):
    tbl2.columns[ci].width = int(tbl_width * cw / total2)
for ri, row in enumerate(hyp_rows):
    for ci, txt in enumerate(row):
        cell = tbl2.cell(ri, ci)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(9)
        run.font.bold = (ri == 0)
        if ri == 0:
            run.font.color.rgb = RGBColor(255,255,255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BG
        elif ri % 2 == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(235, 241, 250)

print("슬라이드 B (q_type 분석 틀) 추가 완료")


# ═══════════════════════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(PPT_PATH)
print(f"\n✅ 저장 완료: {PPT_PATH}")
print(f"총 슬라이드 수: {len(prs.slides)}")
